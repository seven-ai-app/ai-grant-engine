import logging
from pathlib import Path

logger = logging.getLogger(__name__)


async def parse_uploaded_files(file_paths: list[str]) -> list[dict]:
    """Parse uploaded files (PDF, PPTX, DOCX) and extract text content."""
    results = []
    for path_str in file_paths:
        path = Path(path_str)
        if not path.exists():
            logger.warning(f"File not found: {path}")
            continue

        try:
            content = _extract_content(path)
            results.append({
                "filename": path.name,
                "path": str(path),
                "type": path.suffix.lower(),
                "content": content,
            })
        except Exception as e:
            logger.error(f"Error parsing {path}: {e}")
            results.append({
                "filename": path.name,
                "path": str(path),
                "type": path.suffix.lower(),
                "content": "",
                "error": str(e),
            })

    return results


def _extract_content(path: Path) -> str:
    """Extract text from a file based on its extension."""
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(path)
    elif suffix == ".pptx":
        return _extract_pptx(path)
    elif suffix == ".docx":
        return _extract_docx(path)
    elif suffix in (".txt", ".md"):
        return path.read_text(encoding="utf-8")
    else:
        logger.warning(f"Unsupported file type: {suffix}")
        return ""


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF."""
    from PyPDF2 import PdfReader

    reader = PdfReader(str(path))
    text_parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            text_parts.append(text)
    return "\n\n".join(text_parts)


def _extract_pptx(path: Path) -> str:
    """Extract text from PowerPoint presentation."""
    from pptx import Presentation

    prs = Presentation(str(path))
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"--- Slide {slide_num} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
        text_parts.append("\n".join(slide_text))
    return "\n\n".join(text_parts)


def _extract_docx(path: Path) -> str:
    """Extract text from Word document."""
    from docx import Document

    doc = Document(str(path))
    text_parts = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            text_parts.append(text)
    return "\n\n".join(text_parts)
